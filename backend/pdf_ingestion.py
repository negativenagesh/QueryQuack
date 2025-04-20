import pdfplumber
import docx
import pytesseract
from PIL import Image
import streamlit as st
import io
import os
import pandas as pd
import csv
import markdown
import re
from pptx import Presentation
import base64
from io import StringIO, BytesIO
import zipfile
import tempfile

def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file, with optional OCR for scanned PDFs.
    
    Args:
        pdf_file: Either a file-like object or a path to a PDF file.
    """
    text = ""
    tesseract_available = False
    
    # Check if tesseract is available
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        tesseract_available = True
    except (ImportError, pytesseract.TesseractNotFoundError):
        pass
    
    # Determine filename for metadata and error messages
    if isinstance(pdf_file, str):
        filename = os.path.basename(pdf_file)
    else:
        filename = getattr(pdf_file, 'name', 'unknown.pdf')
    
    metadata = {
        "filename": filename,
        "type": "pdf"
    }
    
    try:
        with pdfplumber.open(pdf_file) as pdf:
            metadata["pages"] = len(pdf.pages)
            
            for page_num, page in enumerate(pdf.pages, 1):
                # Try direct text extraction
                extracted = page.extract_text()
                if extracted and extracted.strip():
                    text += f"Page {page_num}: {extracted}\n\n"
                elif tesseract_available:
                    # Fallback to OCR only if tesseract is available
                    try:
                        img = page.to_image().original
                        ocr_text = pytesseract.image_to_string(img)
                        if ocr_text.strip():
                            text += f"Page {page_num} (OCR): {ocr_text}\n\n"
                        else:
                            text += f"Page {page_num}: [No text content detected]\n\n"
                    except Exception as ocr_err:
                        text += f"Page {page_num}: [OCR processing failed]\n\n"
                        st.warning(f"OCR failed on page {page_num} of {filename}: {str(ocr_err)}")
                else:
                    # No OCR available, skip this page
                    text += f"Page {page_num}: [Text extraction failed and OCR not available]\n\n"
            
            # Add table extraction for more comprehensive results
            for page_num, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                if tables:
                    text += f"Tables from Page {page_num}:\n"
                    for i, table in enumerate(tables):
                        text += f"Table {i+1}:\n"
                        for row in table:
                            text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                        text += "\n"
                        
        return text, metadata
    except Exception as e:
        st.error(f"Error extracting text from {filename}: {str(e)}")
        return f"[Error processing file: {filename}]", {"filename": filename, "error": str(e), "type": "pdf"}

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file."""
    try:
        doc = docx.Document(docx_file)
        
        # Extract paragraphs
        paragraphs_text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        
        # Extract tables
        tables_text = ""
        for i, table in enumerate(doc.tables):
            tables_text += f"Table {i+1}:\n"
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                tables_text += row_text + "\n"
            tables_text += "\n"
            
        text = paragraphs_text + "\n\n" + tables_text if tables_text else paragraphs_text
        
        return text, {
            "filename": docx_file.name, 
            "type": "docx", 
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables)
        }
    except Exception as e:
        st.error(f"Error extracting text from {docx_file.name}: {str(e)}")
        return f"[Error processing file: {docx_file.name}]", {"filename": docx_file.name, "error": str(e), "type": "docx"}

def extract_text_from_txt(txt_file):
    """Extract text from a plain text file."""
    try:
        text = txt_file.getvalue().decode('utf-8')
        return text, {"filename": txt_file.name, "type": "txt", "size": len(text)}
    except UnicodeDecodeError:
        # Try different encodings if UTF-8 fails
        try:
            text = txt_file.getvalue().decode('latin-1')
            return text, {"filename": txt_file.name, "type": "txt", "size": len(text), "encoding": "latin-1"}
        except Exception as e:
            st.error(f"Error extracting text from {txt_file.name}: {str(e)}")
            return f"[Error processing file: {txt_file.name}]", {"filename": txt_file.name, "error": str(e), "type": "txt"}
    except Exception as e:
        st.error(f"Error extracting text from {txt_file.name}: {str(e)}")
        return f"[Error processing file: {txt_file.name}]", {"filename": txt_file.name, "error": str(e), "type": "txt"}

def extract_text_from_csv(csv_file):
    """Extract text from a CSV file."""
    try:
        # Read CSV data
        content = csv_file.getvalue().decode('utf-8')
        csv_data = StringIO(content)
        reader = csv.reader(csv_data)
        rows = list(reader)
        
        # Format as text
        text = ""
        for row in rows:
            text += " | ".join(row) + "\n"
        
        return text, {
            "filename": csv_file.name, 
            "type": "csv", 
            "rows": len(rows),
            "columns": len(rows[0]) if rows else 0
        }
    except Exception as e:
        st.error(f"Error extracting text from {csv_file.name}: {str(e)}")
        return f"[Error processing file: {csv_file.name}]", {"filename": csv_file.name, "error": str(e), "type": "csv"}

def extract_text_from_excel(excel_file):
    """Extract text from Excel files."""
    try:
        # Read the Excel file
        df_dict = pd.read_excel(excel_file, sheet_name=None)
        
        text = ""
        sheet_count = 0
        total_rows = 0
        
        # Process each sheet
        for sheet_name, df in df_dict.items():
            sheet_count += 1
            rows = len(df)
            total_rows += rows
            
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
        
        return text, {
            "filename": excel_file.name, 
            "type": "excel", 
            "sheets": sheet_count,
            "total_rows": total_rows
        }
    except Exception as e:
        st.error(f"Error extracting text from {excel_file.name}: {str(e)}")
        return f"[Error processing file: {excel_file.name}]", {"filename": excel_file.name, "error": str(e), "type": "excel"}

def extract_text_from_image(image_file):
    """Extract text from images using OCR."""
    try:
        # Open image with PIL
        image = Image.open(image_file)
        
        # Perform OCR
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            text = "[No text detected in image]"
        
        return text, {
            "filename": image_file.name, 
            "type": "image",
            "dimensions": f"{image.width}x{image.height}"
        }
    except Exception as e:
        st.error(f"Error extracting text from {image_file.name}: {str(e)}")
        return f"[Error processing file: {image_file.name}]", {"filename": image_file.name, "error": str(e), "type": "image"}

def extract_text_from_pptx(pptx_file):
    """Extract text from PowerPoint files."""
    try:
        prs = Presentation(pptx_file)
        text = ""
        
        for i, slide in enumerate(prs.slides, 1):
            text += f"Slide {i}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
            text += "\n"
        
        if not text.strip():
            text = "[No text content found in presentation]"
            
        return text, {
            "filename": pptx_file.name,
            "type": "pptx",
            "slides": len(prs.slides)
        }
    except Exception as e:
        st.error(f"Error extracting text from {pptx_file.name}: {str(e)}")
        return f"[Error processing file: {pptx_file.name}]", {"filename": pptx_file.name, "error": str(e), "type": "pptx"}

def extract_text_from_html(html_file):
    """Extract text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        
        content = html_file.getvalue().decode('utf-8')
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.extract()
            
        # Get text
        text = soup.get_text()
        
        # Break into lines and remove leading/trailing space
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Remove blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text, {
            "filename": html_file.name,
            "type": "html",
            "size": len(content)
        }
    except Exception as e:
        st.error(f"Error extracting text from {html_file.name}: {str(e)}")
        return f"[Error processing file: {html_file.name}]", {"filename": html_file.name, "error": str(e), "type": "html"}

def extract_text_from_markdown(md_file):
    """Extract text from Markdown files."""
    try:
        content = md_file.getvalue().decode('utf-8')
        # Convert markdown to plain text by removing markup
        # Remove headers
        text = re.sub(r'#{1,6}\s+', '', content)
        # Remove bold/italic
        text = re.sub(r'\*\*?(.*?)\*\*?', r'\1', text)
        # Remove links
        text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
        # Remove images
        text = re.sub(r'!\[(.*?)\]\(.*?\)', r'\1', text)
        # Remove blockquotes
        text = re.sub(r'^\s*>\s+', '', text, flags=re.MULTILINE)
        # Remove code blocks
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        # Remove inline code
        text = re.sub(r'`(.*?)`', r'\1', text)
        
        return text, {
            "filename": md_file.name,
            "type": "markdown",
            "size": len(content)
        }
    except Exception as e:
        st.error(f"Error extracting text from {md_file.name}: {str(e)}")
        return f"[Error processing file: {md_file.name}]", {"filename": md_file.name, "error": str(e), "type": "markdown"}

def extract_text_from_zip(zip_file):
    """Extract text from zip files by processing each contained file."""
    try:
        # Create a temporary directory to extract files
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create a BytesIO object from the zip file
            zip_data = BytesIO(zip_file.getvalue())
            
            # Extract all files to the temporary directory
            with zipfile.ZipFile(zip_data) as z:
                z.extractall(temp_dir)
            
            # Process each file in the zip
            extracted_text = f"Contents of ZIP file: {zip_file.name}\n\n"
            file_count = 0
            
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    relative_path = os.path.relpath(file_path, temp_dir)
                    file_count += 1
                    
                    # Try to determine file type and extract if possible
                    extracted_text += f"File: {relative_path}\n"
                    extracted_text += "-" * 50 + "\n"
                    
                    try:
                        # Extract based on file extension
                        if file.lower().endswith('.txt'):
                            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                                extracted_text += f.read()[:1000] + "... [truncated]\n\n"
                        elif file.lower().endswith(('.jpg', '.jpeg', '.png')):
                            extracted_text += "[Image file - OCR not performed on zip contents]\n\n"
                        elif file.lower().endswith('.pdf'):
                            extracted_text += "[PDF file - not processed within zip]\n\n"
                        else:
                            extracted_text += "[File type not processed within zip]\n\n"
                    except Exception as inner_e:
                        extracted_text += f"[Error processing this file: {str(inner_e)}]\n\n"
            
            return extracted_text, {
                "filename": zip_file.name,
                "type": "zip",
                "files": file_count
            }
    except Exception as e:
        st.error(f"Error extracting text from {zip_file.name}: {str(e)}")
        return f"[Error processing file: {zip_file.name}]", {"filename": zip_file.name, "error": str(e), "type": "zip"}

def ingest_files(uploaded_files):
    """Process multiple uploaded files of various types."""
    if not uploaded_files:
        return [], []
    
    results = []
    metadata_list = []
    
    # Create a progress bar
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    total_files = len(uploaded_files)
    
    for i, uploaded_file in enumerate(uploaded_files):
        filename = uploaded_file.name
        status_text.text(f"Processing file {i+1}/{total_files}: {filename}")
        
        try:
            # Determine file type by extension and call appropriate handler
            if filename.lower().endswith('.pdf'):
                text, metadata = extract_text_from_pdf(uploaded_file)
            elif filename.lower().endswith('.docx'):
                text, metadata = extract_text_from_docx(uploaded_file)
            elif filename.lower().endswith('.txt'):
                text, metadata = extract_text_from_txt(uploaded_file)
            elif filename.lower().endswith(('.csv')):
                text, metadata = extract_text_from_csv(uploaded_file)
            elif filename.lower().endswith(('.xlsx', '.xls')):
                text, metadata = extract_text_from_excel(uploaded_file)
            elif filename.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
                text, metadata = extract_text_from_image(uploaded_file)
            elif filename.lower().endswith(('.pptx', '.ppt')):
                text, metadata = extract_text_from_pptx(uploaded_file)
            elif filename.lower().endswith(('.html', '.htm')):
                text, metadata = extract_text_from_html(uploaded_file)
            elif filename.lower().endswith(('.md', '.markdown')):
                text, metadata = extract_text_from_markdown(uploaded_file)
            elif filename.lower().endswith('.zip'):
                text, metadata = extract_text_from_zip(uploaded_file)
            else:
                # Try to process as text file for unknown extensions
                try:
                    text, metadata = extract_text_from_txt(uploaded_file)
                except:
                    text = f"[Unsupported file type: {filename}]"
                    metadata = {"filename": filename, "error": "Unsupported file type", "type": "unknown"}
            
            # Add file index for tracking
            metadata["file_index"] = i
            
            results.append(text)
            metadata_list.append(metadata)
            
        except Exception as e:
            st.error(f"Error processing file {filename}: {str(e)}")
            results.append(f"[Error processing file: {filename}]")
            metadata_list.append({
                "filename": filename, 
                "error": str(e), 
                "type": "error",
                "file_index": i
            })
        
        # Update progress
        progress = (i + 1) / total_files
        progress_bar.progress(progress)
    
    # Clear progress indicators when done
    progress_bar.empty()
    status_text.empty()
    
    if results:
        st.success(f"Successfully processed {len(results)} files")
    
    return results, metadata_list

# For backward compatibility
def ingest_file(uploaded_file):
    """Process a single uploaded file (for backward compatibility)."""
    results, metadata_list = ingest_files([uploaded_file])
    return results[0] if results else "", metadata_list[0] if metadata_list else {}